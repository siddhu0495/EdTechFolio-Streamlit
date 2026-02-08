import streamlit as st
import openai
import ollama
import io
import pandas as pd
import json
import os
import hashlib
import requests
from bs4 import BeautifulSoup
from youtube_transcript_api import YouTubeTranscriptApi
from googlesearch import search
from groq import Groq
# import google.generativeai as genai
from google import genai
from google.genai import types
from fpdf import FPDF
from docx import Document
from pptx import Presentation
from openai import OpenAI

# --- PAGE CONFIG ---
st.set_page_config(
    page_title="EdTech Folio", 
    page_icon="🎓",
    layout="wide",
    initial_sidebar_state="expanded"
)

# --- CUSTOM CSS FOR ENHANCED UI ---
def inject_custom_css():
    st.markdown("""
    <style>
    /* Global Styles */
    .main {
        background-color: #ffffff;
    }
    
    /* Enhanced Service Cards */
    .service-card {
        background: #ffffff;
        border-radius: 12px;
        padding: 20px;
        color: #333333;
        text-align: center;
        height: 200px;
        transition: all 0.3s ease;
        box-shadow: 0 2px 8px rgba(0, 0, 0, 0.08);
        border: 1px solid #e1e5e9;
        position: relative;
        overflow: hidden;
        cursor: pointer;
    }
    
    .service-card:hover {
        transform: translateY(-3px);
        box-shadow: 0 4px 12px rgba(0, 0, 0, 0.12);
        border-color: #4f46e5;
    }
    
    .service-card-content {
        display: flex;
        flex-direction: column;
        justify-content: center;
        align-items: center;
        height: 100%;
    }
    
    .service-icon {
        font-size: 2.2rem;
        margin-bottom: 10px;
        color: #4f46e5;
    }
    
    .service-title {
        font-size: 1.1rem;
        font-weight: 600;
        margin-bottom: 6px;
        color: #1f2937;
    }
    
    .service-desc {
        font-size: 0.9rem;
        color: #6b7280;
        line-height: 1.4;
    }
    
    
    /* Enhanced Button Styles */
    .stButton > button {
        border-radius: 8px !important;
        border: 1px solid #e5e7eb !important;
        background-color: #ffffff !important;
        color: #374151 !important;
        font-weight: 500 !important;
        transition: all 0.2s ease !important;
        box-shadow: 0 1px 3px rgba(0, 0, 0, 0.1) !important;
    }
    
    .stButton > button:hover {
        transform: translateY(-1px) !important;
        box-shadow: 0 2px 6px rgba(0, 0, 0, 0.15) !important;
        border-color: #4f46e5 !important;
        color: #4f46e5 !important;
    }
    
    /* Enhanced Input Fields */
    .stTextInput > div > div > input,
    .stTextArea > div > div > textarea,
    .stNumberInput > div > div > input {
        border-radius: 8px !important;
        border: 1px solid #e5e7eb !important;
        padding: 10px 12px !important;
        font-size: 14px !important;
        transition: all 0.2s ease !important;
        background-color: #ffffff !important;
    }
    
    .stTextInput > div > div > input:focus,
    .stTextArea > div > div > textarea:focus,
    .stNumberInput > div > div > input:focus {
        border-color: #4f46e5 !important;
        box-shadow: 0 0 0 3px rgba(79, 70, 229, 0.1) !important;
        outline: none !important;
    }
    
    /* Enhanced Select Boxes */
    .stSelectbox > div > div > select {
        border-radius: 8px !important;
        border: 1px solid #e5e7eb !important;
        padding: 10px 12px !important;
        font-size: 14px !important;
        background-color: #f9fafb !important;
        transition: all 0.2s ease !important;
    }
    
    .stSelectbox > div > div > select:focus {
        border-color: #4f46e5 !important;
        box-shadow: 0 0 0 3px rgba(79, 70, 229, 0.1) !important;
        outline: none !important;
        background-color: #ffffff !important;
    }
    
    /* Enhanced Sliders */
    .stSlider > div > div > div > div {
        background: #4f46e5 !important;
    }
    
    /* Enhanced Sidebar */
    section[data-testid="stSidebar"] {
        background: #ffffff !important;
        border-right: 1px solid #e5e7eb !important;
    }
    
    section[data-testid="stSidebar"] .css-1d391kg {
        color: #374151 !important;
    }
    
    /* Enhanced Headers */
    h1, h2, h3 {
        color: #111827 !important;
        font-weight: 700 !important;
    }
    
    /* Enhanced Success Messages */
    .stSuccess {
        background-color: #ecfdf5 !important;
        border-color: #bbf7d0 !important;
        color: #065f46 !important;
        border-radius: 8px !important;
        border: 1px solid #bbf7d0 !important;
    }
    
    /* Enhanced Warning Messages */
    .stWarning {
        background-color: #fffbeb !important;
        border-color: #fed7aa !important;
        color: #92400e !important;
        border-radius: 8px !important;
        border: 1px solid #fed7aa !important;
    }
    
    /* Enhanced Error Messages */
    .stError {
        background-color: #fef2f2 !important;
        border-color: #fecaca !important;
        color: #991b1b !important;
        border-radius: 8px !important;
        border: 1px solid #fecaca !important;
    }
    
    /* Enhanced Info Messages */
    .stInfo {
        background-color: #eff6ff !important;
        border-color: #bfdbfe !important;
        color: #1e40af !important;
        border-radius: 8px !important;
        border: 1px solid #bfdbfe !important;
    }
    
    /* Enhanced Expander */
    .streamlit-expanderHeader {
        background-color: #f9fafb !important;
        border-radius: 8px !important;
        border: 1px solid #e5e7eb !important;
        color: #374151 !important;
    }
    
    /* Enhanced Download Button */
    .stDownloadButton > button {
        background: #10b981 !important;
        border: none !important;
        color: white !important;
        border-radius: 8px !important;
        font-weight: 600 !important;
        padding: 10px 20px !important;
        transition: all 0.2s ease !important;
    }
    
    .stDownloadButton > button:hover {
        background: #059669 !important;
        transform: translateY(-1px) !important;
        box-shadow: 0 4px 6px rgba(16, 185, 129, 0.3) !important;
    }
    
    /* Enhanced Navigation Radio */
    .stRadio > div {
        background-color: #f9fafb !important;
        padding: 8px !important;
        border-radius: 8px !important;
        border: 1px solid #e5e7eb !important;
    }
    
    /* Loading Spinner Enhancement */
    .stSpinner > div > div {
        border-top-color: #4f46e5 !important;
        border-left-color: #4f46e5 !important;
    }
    
    /* Content Area Enhancement */
    .generated-content {
        background: #ffffff;
        padding: 20px;
        border-radius: 12px;
        box-shadow: 0 1px 3px rgba(0, 0, 0, 0.1);
        border: 1px solid #e5e7eb;
        margin: 16px 0;
    }
    
    /* Category Tags */
    .category-tag {
        display: inline-block;
        background: #4f46e5;
        color: white;
        padding: 4px 8px;
        border-radius: 8px;
        font-size: 0.75rem;
        font-weight: 600;
        margin-bottom: 8px;
        text-transform: uppercase;
        letter-spacing: 0.5px;
    }
    
    /* Tool Header Enhancement */
    .tool-header {
        color: #111827;
        font-weight: 700;
        margin-bottom: 8px;
    }
    
    /* Responsive Design */
    @media (max-width: 768px) {
        .service-card {
            height: auto !important;
            min-height: 160px !important;
        }
        
        .service-title {
            font-size: 1rem !important;
        }
        
        .service-desc {
            font-size: 0.85rem !important;
        }
        
        .service-icon {
            font-size: 2rem !important;
        }
    }
    </style>
    """, unsafe_allow_html=True)

# Inject CSS at the beginning of the app
inject_custom_css()

# --- LLM INTEGRATION PROVISION ---
def get_llm_response(tool_name, inputs):
    """
    Generates a response using the selected LLM provider (Mock, OpenAI, or Gemini).
    """
    provider = st.session_state.get('api_provider', 'Mock')
    api_key = st.session_state.get('api_key', '')
    ollama_model = st.session_state.get('ollama_model', 'openai/gpt-oss-20b')

    # ollama_base_url = ""
    # try:
    #     # Read from Streamlit secrets if available
    #     if "ollama_base_url" in st.secrets:
    #         ollama_base_url = st.secrets["ollama_base_url"]
    # except Exception:
    #     pass # Fails gracefully if secrets file doesn't exist or key is missing
    truncation_msg = ""

    # Helper to yield string as stream (for Mock or non-streaming fallbacks)
    def stream_text(text):
        yield text

    # 1. Construct the Prompt based on the tool
    prompt = ""
    if tool_name == "YouTube Generator":
        prompt = f"Generate guiding questions for a YouTube video with this URL: {inputs.get('url')}. If you cannot access the video, provide general questions based on the topic inferred from the URL."
    elif tool_name == "Text Question":
        prompt = f"Generate text-dependent questions based on the following text:\n\n{inputs.get('text')}"
    elif tool_name == "Worksheet Generator":
        prompt = f"Create a comprehensive worksheet for the topic: {inputs.get('topic')}."
    elif tool_name == "MCQ Generator":
        prompt = f"Generate {inputs.get('num_questions')} multiple-choice questions based on this context:\n\n{inputs.get('context')}"
    elif tool_name == "Text Summarizer":
        prompt = f"Summarize the following text. Length: {inputs.get('length')}.\n\nText: {inputs.get('text')}"
    elif tool_name == "Text Rewriter":
        prompt = f"Rewrite the following text based on these criteria: {inputs.get('criteria')}.\n\nText: {inputs.get('text')}"
    elif tool_name == "Proof Read":
        prompt = f"Proofread the following text for grammar, spelling, punctuation, and clarity:\n\n{inputs.get('text')}"
    elif tool_name == "Lesson Plan":
        prompt = f"Create a lesson plan for '{inputs.get('topic')}' with a duration of {inputs.get('duration')} minutes."
    elif tool_name == "Report Card":
        prompt = f"Write report card comments for student '{inputs.get('name')}'. Strengths: {inputs.get('strengths')}. Areas for Growth: {inputs.get('growth_areas')}."
    elif tool_name == "Essay Grader":
        prompt = f"Grade the following essay based on this rubric: {inputs.get('rubric')}.\n\nEssay:\n{inputs.get('essay')}"
    elif tool_name == "PPT Generator":
        prompt = f"Create a presentation outline with {inputs.get('num_slides')} slides for the topic: {inputs.get('topic')}."
    elif tool_name == "Question Paper Creation":
        prompt = f"Create a {inputs.get('exam_type')} question paper for '{inputs.get('topic')}' with {inputs.get('num_questions')} questions."
    elif tool_name == "Paper Correction":
        prompt = f"Provide a correction summary for the answer sheet file named '{inputs.get('file_name')}'. (Note: Real file analysis requires OCR integration)."
    elif tool_name == "YouTube or URL Summary":
        action = inputs.get('action')
        content = inputs.get('content')
        url = inputs.get('url')
        if len(content) > 25000:
            truncation_msg = f"\n\n⚠️ **Note:** The text is truncated and complete info is not displayed. Reference: {url}"
            content = content[:25000]
        
        if action == "Complete Text":
            return content + truncation_msg
        elif action == "Translation":
            prompt = f"Translate the following content into English. Provide key insights if applicable.\n\nContent: {content}"
        elif "Summary" in action:
            words_instruction = f" in approximately {action.split('(')[1].split(' ')[0]} words" if "(" in action else ""
            prompt = f"Summarize the following content{words_instruction}. Provide key insights.\n\nContent: {content}"
    elif tool_name == "Google Search Info":
        prompt = f"Synthesize the following search results for the topic '{inputs.get('topic')}'. Use the 'Content Snippet' if available; otherwise, rely on the 'Description'. You MUST cite the URLs for every claim. If you cannot find specific details, list the URLs as resources.\n\nSearch Data:\n{inputs.get('search_results')}"
    elif tool_name == "Google News":
        prompt = f"Summarize the key news stories regarding '{inputs.get('topic')}' based on the following headlines. Highlight the most important events.\n\nNews Data:\n{inputs.get('news_results')}"

    # 2. Handle Real API Calls
    if provider != "Mock":
        if provider != "Ollama" and not api_key:
            return "⚠️ Please enter an API Key in the sidebar settings."
        
        try:
            if provider == "OpenAI":
                client = openai.OpenAI(api_key=api_key)
                response = client.chat.completions.create(
                    model="gpt-3.5-turbo",
                    messages=[{"role": "user", "content": prompt}]
                )
                return stream_text(response.choices[0].message.content)
            elif provider == "Gemini":
                client = genai.Client(api_key=api_key)
                response = client.models.generate_content(
                    model='gemini-2.0-flash', contents=prompt
                )
                return stream_text(response.text)
            elif provider == "Ollama":
                # --- GROQ INTEGRATION (Replaces Ollama logic) ---
                # groq_api_key = st.secrets.get("GROQ_API_KEY", "")
                groq_api_key = st.secrets["GROQ_API_KEY"]
                if not groq_api_key:
                    return "⚠️ GROQ_API_KEY not found in secrets.toml."
                
                # gpt-oss:20b
                # client = OpenAI(
                #     base_url="https://api.groq.com/openai/v1",
                #     # api_key=st.secrets["GROQ_API_KEY"]
                #     api_key=groq_api_key
                # )

                # client = OpenAI(
                #     base_url="http://localhost:11434/v1",
                #     # api_key=st.secrets["GROQ_API_KEY"]
                #     api_key="ollama"
                # )
                
                client = Groq(api_key=groq_api_key)
                # chat_completion = client.chat.completions.create(
                stream = client.chat.completions.create(
                    messages=[{"role": "user", "content": prompt}],
                    model=ollama_model,
                    stream=True,
                )
                
                def generate_stream():
                    for chunk in stream:
                        content = chunk.choices[0].delta.content
                        if content: yield content
                    if tool_name == "YouTube or URL Summary" and truncation_msg:
                        yield truncation_msg
                
                return generate_stream()
                
                # --- DISABLED OLLAMA/NGROK LOGIC (Kept for reference) ---
                # if ollama_base_url: ... client = ollama.Client(...) ...
                # else: ... response = ollama.chat(...) ...
            
                ### This code is commented. Can be removed later - Start
                # # Ensure Ollama is running locally (default port 11434) or via ngrok
                # if ollama_base_url:
                #     client = ollama.Client(host=ollama_base_url, headers={'ngrok-skip-browser-warning': 'true'})
                #     response = client.chat(model=ollama_model, messages=[{'role': 'user', 'content': prompt}])
                # else:
                #     response = ollama.chat(model=ollama_model, messages=[{'role': 'user', 'content': prompt}])
                # response_text = response['message']['content']
                ### This code is commented. Can be removed later - End

            # if tool_name == "YouTube or URL Summary" and truncation_msg:
            #     return response_text + truncation_msg
            # return response_text
        except Exception as e:
            if provider == "Ollama" and "404" in str(e):
                return stream_text(f"⚠️ Model '{ollama_model}' not found locally. Please run `ollama pull {ollama_model}` in your terminal.")
            return stream_text(f"❌ Error: {str(e)}")
    
    # 3. Fallback to Mock Logic
    if tool_name == "YouTube Generator":
        return stream_text("Sample Output: What are the three key points mentioned about photosynthesis in the video?")
    elif tool_name == "Text Question":
        return stream_text("Sample Output: Based on the text, why did the character decide to leave?")
    elif tool_name == "Worksheet Generator":
        return stream_text("Sample Worksheet:\n\n**The Water Cycle**\n\n1. Fill in the blank: The process of water turning into vapor is called ______.\n2. What is precipitation?")
    elif tool_name == "MCQ Generator":
        return stream_text(f"Generated {inputs.get('num_questions')} questions.\n\nSample Question: What is the capital of France?\nA) Berlin B) Madrid C) Paris D) Rome")
    elif tool_name == "Text Summarizer":
        return stream_text(f"This is a {inputs.get('length', '').lower()} summary of the provided text, focusing on the key points.")
    elif tool_name == "Text Rewriter":
        return stream_text(f"Here is the rewritten text, adapted to be '{inputs.get('criteria')}'. ")
    elif tool_name == "Proof Read":
        return stream_text("Sample Output: The original text had a few grammatical errors which have been corrected for clarity and flow.")
    elif tool_name == "Lesson Plan":
        return stream_text(f"**Lesson Plan: {inputs.get('topic')} ({inputs.get('duration')} mins)**\n\n*   **Objective:** Students will be able to...\n*   **Activities:** Warm-up, Direct Instruction, Guided Practice...")
    elif tool_name == "Report Card":
        return stream_text(f"{inputs.get('name')} is a pleasure to have in class. He/She consistently demonstrates {inputs.get('strengths')}. To further improve, {inputs.get('name')} should focus on {inputs.get('growth_areas')}.")
    elif tool_name == "Essay Grader":
        return stream_text("**Grade: 4/5**\n\n**Feedback:** The essay presents a strong argument but could be improved by providing more specific examples. There are minor grammatical errors.")
    elif tool_name == "PPT Generator":
        return stream_text(f"**{inputs.get('topic')} - {inputs.get('num_slides')} Slide Outline**\n\n1.  **Title Slide:** {inputs.get('topic')}\n2.  **Introduction:** What is AI?\n3.  **History:** Early Concepts\n...")
    elif tool_name == "Question Paper Creation":
        return stream_text(f"**{inputs.get('exam_type')} Exam - {inputs.get('topic')}**\n\n**Total Questions: {inputs.get('num_questions')}**\n\n1.  Discuss the primary causes of {inputs.get('topic')}.\n2.  Define the term 'Blitzkrieg'.")
    elif tool_name == "Paper Correction":
        return stream_text("**Correction Summary:**\n- Question 1: Correct\n- Question 2: Partially incorrect. Key detail missed.\n\n**Marks: 8/10**")
    elif tool_name == "YouTube or URL Summary":
        return stream_text(f"**{inputs.get('action')} Output:**\n\nThis is a simulated output for the content from {inputs.get('url')}.")
    elif tool_name == "Google Search Info":
        return stream_text(f"**Latest Info on {inputs.get('topic')}:**\n\n(Mock Summary) The search results indicate significant interest in {inputs.get('topic')}. Key developments include...")
    elif tool_name == "Google News":
        return stream_text(f"**News Update for {inputs.get('topic')}:**\n\n(Mock News) Several outlets are reporting on {inputs.get('topic')}. Major headlines involve...")
    
    return stream_text("Output generated.")

# --- SCRAPING UTILS ---
def get_video_id(url):
    """Extracts video ID from YouTube URL."""
    if "youtu.be" in url:
        return url.split("/")[-1].split("?")[0]
    if "v=" in url:
        return url.split("v=")[1].split("&")[0]
    return None

def fetch_url_content(url):
    """Fetches text content from a URL (YouTube transcript or Webpage text)."""
    try:
        if "youtube.com" in url or "youtu.be" in url:
            video_id = get_video_id(url)
            # print(video_id,'video_id')
            if not video_id: return "Error: Invalid YouTube URL."
            try:
                # transcript_list = YouTubeTranscriptApi.get_transcript(video_id)
                # transcriptYouTube = YouTubeTranscriptApi()
                # transcript_list = transcriptYouTube.fetch(video_id)
                transcript_list = YouTubeTranscriptApi().fetch(video_id)   # This will fetch the transcript for the video
                # transcript_list = YouTubeTranscriptApi().list(video_id)    # This will list all the available languages (transcripts) for the video
       
            except AttributeError:
                return "Error: Library conflict detected. Please ensure you do not have a file named 'youtube_transcript_api.py' in your project folder."
            # print('transcript_list start',transcript_list,'transcript_list end')
            transcript_text = " ".join([t.text for t in transcript_list])
            # print('transcript_text start',transcript_text,'transcript_text end')
            return transcript_text
        else:
            headers = {
                'User-Agent': 'Mozilla/5.0',
                'ngrok-skip-browser-warning': 'true'
            }
            response = requests.get(url, headers=headers, timeout=10)
            if response.status_code != 200: return f"Error: Unable to fetch URL (Status: {response.status_code})"
            soup = BeautifulSoup(response.text, 'html.parser')
            for script in soup(["script", "style"]): script.extract()
            return soup.get_text(separator=' ', strip=True)
    except Exception as e:
        return f"Error reading content: {str(e)}"

def fetch_google_search(query):
    """Fetches Google Search results."""
    try:
        # Try advanced search first (returns objects with title/desc)
        try:
            results = list(search(query, num_results=10, advanced=True))
        except:
            results = []

        # Fallback to basic search if advanced fails or returns empty
        if not results:
            try:
                # basic search returns strings (URLs)
                urls = list(search(query, num_results=10))
                # Create mock objects for consistency
                class MockResult:
                    def __init__(self, url):
                        self.url = url
                        self.title = "Title unavailable"
                        self.description = "Description unavailable"
                results = [MockResult(u) for u in urls]
            except:
                return "No search results found."

        output = []
        for i, r in enumerate(results):
            # Attempt to scrape the content of the page
            scraped_content = fetch_url_content(r.url)
            # Truncate scraped content to avoid context window overflow (approx 1000 chars per source)
            if "Error" in scraped_content:
                # If scraping fails, just note it so LLM knows to use description
                scraped_content = "(Content could not be scraped. Use description.)"
            else:
                scraped_content = scraped_content[:1000].replace("\n", " ") + "..."
            
            output.append(f"Source {i+1}:\nTitle: {r.title}\nURL: {r.url}\nDescription: {r.description}\nContent Snippet: {scraped_content}\n")
            
        return "\n".join(output)
    except Exception as e:
        return f"Error fetching search results: {str(e)}"

def fetch_google_news(query):
    """Fetches Google News RSS feed."""
    try:
        url = f"https://news.google.com/rss/search?q={query}&hl=en-US&gl=US&ceid=US:en"
        response = requests.get(url, timeout=5)
        # Use html.parser to handle XML tags in a forgiving way
        soup = BeautifulSoup(response.content, 'html.parser')
        items = soup.find_all('item', limit=5)
        output = []
        for item in items:
            title = item.title.text if item.title else "No Title"
            pub = item.pubdate.text if item.find('pubdate') else ""
            output.append(f"Headline: {title}\nDate: {pub}")
        return "\n\n".join(output)
    except Exception as e:
        return f"Error fetching news: {str(e)}"

# --- FILE GENERATION UTILS ---
def generate_file_content(content, file_format):
    """Converts text content to the specified file format."""
    if file_format == "txt":
        return content.encode("utf-8"), "text/plain"
    
    elif file_format == "html":
        html_content = f"<html><body><pre>{content}</pre></body></html>"
        return html_content.encode("utf-8"), "text/html"
    
    elif file_format == "pdf":
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Arial", size=12)
        # FPDF doesn't handle unicode well by default, so we replace/encode
        # For production, use a unicode font.
        safe_content = content.encode('latin-1', 'replace').decode('latin-1')
        pdf.multi_cell(0, 10, safe_content)
        return pdf.output(dest='S').encode('latin-1'), "application/pdf"
    
    elif file_format == "docx":
        doc = Document()
        doc.add_paragraph(content)
        bio = io.BytesIO()
        doc.save(bio)
        return bio.getvalue(), "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    
    elif file_format == "ppt":
        prs = Presentation()
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Generated Content"
        slide.placeholders[1].text = content
        bio = io.BytesIO()
        prs.save(bio)
        return bio.getvalue(), "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    
    elif file_format == "csv":
        df = pd.DataFrame([{"Content": content}])
        return df.to_csv(index=False).encode('utf-8'), "text/csv"
    
    elif file_format == "xlsx":
        df = pd.DataFrame([{"Content": content}])
        bio = io.BytesIO()
        df.to_excel(bio, index=False)
        return bio.getvalue(), "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        
    return None, None

def render_download_section(content, key_suffix):
    """Renders the download options section."""
    st.divider()
    st.subheader("💾 Download Options")
    col1, col2 = st.columns([2, 1])
    with col1:
        file_format = st.selectbox("Select Format", ["txt", "html", "pdf", "docx", "ppt", "csv", "xlsx"], key=f"fmt_{key_suffix}")
    with col2:
        file_data, mime = generate_file_content(content, file_format)
        st.download_button("Download File", data=file_data, file_name=f"generated_content.{file_format}", mime=mime, key=f"btn_{key_suffix}")

# --- DATA PERSISTENCE ---
DATA_FILE = "edtech_data.json"

def load_data():
    if not os.path.exists(DATA_FILE):
        return {"users": {}}
    try:
        with open(DATA_FILE, "r") as f:
            return json.load(f)
    except:
        return {"users": {}}

def save_data(data):
    with open(DATA_FILE, "w") as f:
        json.dump(data, f, indent=4)

def hash_password(password):
    return hashlib.sha256(password.encode()).hexdigest()

def register_user(username, password):
    data = load_data()
    if username in data["users"]:
        return False
    data["users"][username] = {
        "password": hash_password(password),
        "history": []
    }
    save_data(data)
    return True

#admin: password123
#testuser: password
#test: sid

def authenticate_user(username, password):
    if username == "test" and password == "sid":
        return True

    data = load_data()
    if username not in data["users"]:
        return False
    if data["users"][username]["password"] == hash_password(password):
        return True
    return False

def get_user_history(username):
    data = load_data()
    return data["users"].get(username, {}).get("history", [])

def save_user_history(username, entry):
    data = load_data()
    if username in data["users"]:
        data["users"][username]["history"].append(entry)
        save_data(data)

def add_history_entry(tool, input_data, output):
    entry = {"tool": tool, "input": input_data, "output": output}
    st.session_state.history.append(entry)
    if st.session_state.get('username'):
        save_user_history(st.session_state.username, entry)

# --- INITIALIZE SESSION STATE ---
if 'logged_in' not in st.session_state:
    st.session_state.logged_in = False
if 'username' not in st.session_state:
    st.session_state.username = ""
if 'history' not in st.session_state:
    st.session_state.history = []
if 'current_tool' not in st.session_state:
    st.session_state.current_tool = None

def login_page():
    st.container()
    col1, col2, col3 = st.columns([1, 2, 1])
    with col2:
        st.title("🔐 Login to EdTech Hub")
        tab1, tab2 = st.tabs(["Login", "Sign Up"])
        
        with tab1:
            username = st.text_input("Username", value="test", key="login_user")
            password = st.text_input("Password", value="sid", type="password", key="login_pass")
            if st.button("Login", use_container_width=True):
                if authenticate_user(username, password):
                    st.session_state.logged_in = True
                    st.session_state.username = username
                    st.session_state.history = get_user_history(username)
                    st.rerun()
                else:
                    st.error("Invalid username or password")
        
        with tab2:
            new_user = st.text_input("Username", key="signup_user")
            new_pass = st.text_input("Password", type="password", key="signup_pass")
            if st.button("Sign Up", use_container_width=True):
                if new_user and new_pass:
                    if register_user(new_user, new_pass):
                        st.success("Account created! Please login.")
                    else:
                        st.error("Username already exists.")
                else:
                    st.error("Please fill in all fields")

if not st.session_state.logged_in:
    login_page()
else:
    # --- SIDEBAR ---
    with st.sidebar:
        st.title("EdTech Hub")
        st.write(f"👤 **{st.session_state.username}**")
        st.divider()
        
        def reset_tool_state():
            st.session_state.current_tool = None
            st.session_state.last_tool_output = None

        nav_option = st.radio("Navigation", ["Home", "Magic Tools", "Chat History"], on_change=reset_tool_state, key="nav_selection")

        st.divider()
        st.subheader("⚙️ Settings")
        st.session_state.api_provider = st.selectbox("Select Model Provider", ["Ollama"], key="api_provider_select")
        
        if st.session_state.api_provider == "Ollama":
            # Add more models to this list as needed
            # st.session_state.ollama_model = st.selectbox("Select Ollama Model", ["gpt-oss", "deepseek-r1", "llama3.2", "qwen", "qwen-vl", "mistral", "phi", "gemma", "gemma3", "translategemma"], key="ollama_model_select")
            
            ### This code is commented. Can be removed later - Start
            # # Initialize default URL from secrets if available, but allow UI override
            # default_ollama_url = ""
            # try:
            #     default_ollama_url = st.secrets.get("ollama_base_url", "")
            # except:
            #     pass
            
            # if "ollama_base_url" not in st.session_state:
            #     st.session_state.ollama_base_url = default_ollama_url
                
            # st.text_input("Ollama Base URL (for remote/ngrok)", key="ollama_base_url", placeholder="https://your-ngrok-url.ngrok-free.app")
            # st.caption("Leave empty for localhost. Required if deploying on Streamlit Cloud.")
            
            # if st.session_state.ollama_base_url:
            #     if st.button("Test Connection", key="test_ollama_conn"):
            #         try:
            #             # Ollama root endpoint usually returns "Ollama is running"
            #             test_url = st.session_state.ollama_base_url.rstrip('/')
            #             response = requests.get(test_url, headers={'ngrok-skip-browser-warning': 'true'}, timeout=5)
            #             if response.status_code == 200:
            #                 st.success(f"✅ Connected: {response.text}")
            #             else:
            #                 st.error(f"❌ Error {response.status_code}: {response.text}")
            #         except Exception as e:
            #             st.error(f"❌ Connection Failed: {str(e)}")
            ### This code is commented. Can be removed later - End


            # Updated list to Groq-supported models
            st.session_state.ollama_model = st.selectbox("Select Model", [
                "openai/gpt-oss-20b", "llama-3.2-3b-preview", "deepseek-v3", "mistral-saba-24b", 
                "gemma-3-27b-it", "deepseek-r1-distill-llama-70b"
            ], key="ollama_model_select")

            # st.session_state.ollama_model = st.selectbox("Select Model", [
            #     "openai/gpt-oss-20b","gpt-oss-20b", "llama3-8b-8192", "llama3-70b-8192", "mixtral-8x7b-32768", 
            #     "gemma-7b-it", "gemma2-9b-it", "deepseek-r1-distill-llama-70b",
            #     "llama-3.2-1b-preview", "llama-3.2-3b-preview", "llama-3.2-11b-vision-preview"
            # ], key="ollama_model_select")

        elif st.session_state.api_provider != "Mock":
            st.session_state.api_key = st.text_input(f"Enter {st.session_state.api_provider} API Key", type="password", key="api_key_input")
            
        st.divider()
        if st.button("Logout", key="logout_btn"):
            st.session_state.logged_in = False
            st.session_state.username = ""
            st.session_state.history = []
            st.rerun()

    # --- MAIN PANEL ---

    if nav_option == "Home":
        st.title("🏠 Welcome to EdTech Folio")
        st.markdown("""
        ### Your AI-Powered Education Assistant
        
        Select **Magic Tools** from the sidebar to access our suite of tools:
        - 🎬 **YouTube Generator**
        - 📝 **Worksheet & MCQ Generators**
        - 🍎 **Lesson Planning**
        - And much more!
        
        Check **Chat History** to view your saved generations.
        """)
        
        col1, col2 = st.columns(2)
        with col1:
            st.button("✨ Go to Magic Tools", use_container_width=True, on_click=lambda: (st.session_state.update(nav_selection="Magic Tools"), reset_tool_state()))
        with col2:
            st.button("📜 View Chat History", use_container_width=True, on_click=lambda: (st.session_state.update(nav_selection="Chat History"), reset_tool_state()))

    elif nav_option == "Chat History":
        st.title("📜 Chat History")
        if not st.session_state.history:
            st.info("No history available yet. Use the Magic Tools to generate content!")
        else:
            for i, entry in enumerate(reversed(st.session_state.history)):
                with st.expander(f"#{len(st.session_state.history)-i}: {entry['tool']}", expanded=True):
                    st.markdown("**User Input:**")
                    if isinstance(entry.get('input'), dict):
                        for key, value in entry['input'].items():
                            st.markdown(f"- **{key.replace('_', ' ').title()}:** {value}")
                    else:
                        st.write(entry.get('input'))
                    st.markdown("**Model Answer:**")
                    st.write(entry['output'])

    elif nav_option == "Magic Tools":
        # Check if a specific tool is selected
        if st.session_state.current_tool:
            just_generated = False
            if st.button("← Back to Magic Tools"):
                st.session_state.current_tool = None
                st.rerun()
            
            tool = st.session_state.current_tool
            
            # --- TOOL: YouTube Generator ---
            if tool == "YouTube Generator":
                st.title("🎬 YouTube Generator")
                st.write("Generate guiding questions aligned to a YouTube video.")
                url = st.text_input("YouTube Video URL", placeholder="https://www.youtube.com/watch?v=...")
                if st.button("Generate Guiding Questions", key="yt_gen"):
                    if url:
                        with st.spinner("🎬 Analyzing video content and generating questions..."):
                            stream = get_llm_response(tool, {"url": url})
                            output = st.write_stream(stream)
                            st.session_state.last_tool_output = {'tool': tool, 'content': output}
                            add_history_entry("YouTube Generator", {"url": url}, output)
                            just_generated = True
                    else:
                        st.warning("Please enter a YouTube URL.")

            # --- TOOL: Text Question ---
            elif tool == "Text Question":
                st.title("❓ Text Question Generator")
                st.write("Generate text-dependent questions for students based on any text.")
                text = st.text_area("Paste your text here:", height=250)
                if st.button("Generate Questions", key="txt_q_gen"):
                    if text:
                        with st.spinner("❓ Analyzing text and generating thought-provoking questions..."):
                            stream = get_llm_response(tool, {"text": text})
                            output = st.write_stream(stream)
                            st.session_state.last_tool_output = {'tool': tool, 'content': output}
                            add_history_entry("Text Question", text[:100] + "...", output)
                            just_generated = True
                    else:
                        st.warning("Please paste some text.")

            # --- TOOL: Worksheet Generator ---
            elif tool == "Worksheet Generator":
                st.title("📝 Worksheet Generator")
                st.write("Generate a Worksheet based on any topic or text.")
                topic = st.text_input("Topic or Text", placeholder="e.g., The Water Cycle")
                if st.button("Generate Worksheet", key="ws_gen"):
                    if topic:
                        with st.spinner("📝 Creating comprehensive worksheet with exercises and activities..."):
                            stream = get_llm_response(tool, {"topic": topic})
                            output = st.write_stream(stream)
                            st.session_state.last_tool_output = {'tool': tool, 'content': output}
                            add_history_entry("Worksheet Generator", {"topic": topic}, output)
                            just_generated = True
                    else:
                        st.warning("Please enter a topic.")

            # --- TOOL: MCQ Generator ---
            elif tool == "MCQ Generator":
                st.title("📝 MCQ Generator")
                st.write("Create a multiple choice assessment based on any topic, standard(s) or criteria!")
                context = st.text_area("Paste your source text, topic, or criteria here:", height=200)
                num_questions = st.slider("Number of questions", 1, 10, 5)
                if st.button("Generate Quiz", key="mcq_gen"):
                    if context:
                        with st.spinner("📝 Creating challenging multiple choice questions with answer keys..."):
                            stream = get_llm_response(tool, {"context": context, "num_questions": num_questions})
                            output = st.write_stream(stream)
                            st.session_state.last_tool_output = {'tool': tool, 'content': output}
                            add_history_entry("MCQ Generator", {"context": context[:100] + "...", "num_questions": num_questions}, output)
                            just_generated = True
                    else:
                        st.warning("Please provide some context.")

            # --- TOOL: Text Summarizer ---
            elif tool == "Text Summarizer":
                st.title("✂️ Text Summarizer")
                st.write("Take any text and summarize it in whatever length you choose.")
                text = st.text_area("Paste the text to summarize:", height=250)
                length = st.select_slider("Summary Length", options=["Short", "Medium", "Long"], value="Medium")
                if st.button("Summarize Text", key="sm_gen"):
                    if text:
                        with st.spinner("✂️ Condensing text to essential points while preserving key information..."):
                            stream = get_llm_response(tool, {"text": text, "length": length})
                            output = st.write_stream(stream)
                            st.session_state.last_tool_output = {'tool': tool, 'content': output}
                            add_history_entry("Text Summarizer", {"text": text[:100] + "...", "length": length}, output)
                            just_generated = True
                    else:
                        st.warning("Please paste some text to summarize.")

            # --- TOOL: Text Rewriter ---
            elif tool == "Text Rewriter":
                st.title("✍️ Text Rewriter")
                st.write("Take any text and rewrite it with custom criteria however you'd like!")
                text = st.text_area("Paste the text to rewrite:", height=200)
                criteria = st.text_input("Rewriting Criteria", placeholder="e.g., Make it more professional, simplify for a 5th grader")
                if st.button("Rewrite Text", key="rw_gen"):
                    if text and criteria:
                        with st.spinner("✍️ Transforming text to match your specified criteria and style..."):
                            stream = get_llm_response(tool, {"text": text, "criteria": criteria})
                            output = st.write_stream(stream)
                            st.session_state.last_tool_output = {'tool': tool, 'content': output}
                            add_history_entry("Text Rewriter", {"text": text[:100] + "...", "criteria": criteria}, output)
                            just_generated = True
                    else:
                        st.warning("Please provide both text and criteria.")

            # --- TOOL: Proof Read ---
            elif tool == "Proof Read":
                st.title("🧐 Proofreader")
                st.write("Have any text proofread for grammar, spelling, punctuation, and clarity.")
                text = st.text_area("Paste the text to proofread:", height=250)
                if st.button("Proofread Text", key="pr_gen"):
                    if text:
                        with st.spinner("🧐 Analyzing text for grammar, spelling, punctuation, and clarity issues..."):
                            stream = get_llm_response(tool, {"text": text})
                            output = st.write_stream(stream)
                            st.session_state.last_tool_output = {'tool': tool, 'content': output}
                            add_history_entry("Proof Read", text[:100] + "...", output)
                            just_generated = True
                    else:
                        st.warning("Please paste some text to proofread.")

            # --- TOOL: Lesson Plan ---
            elif tool == "Lesson Plan":
                st.title("🍎 Lesson Plan Generator")
                st.write("Generate a lesson plan for a topic or objective you're teaching.")
                topic = st.text_input("Topic/Objective", placeholder="e.g., Introduction to Python loops")
                duration = st.number_input("Class Duration (minutes)", min_value=15, max_value=120, value=45)
                if st.button("Generate Lesson Plan", key="lp_gen"):
                    if topic:
                        with st.spinner("🍎 Designing comprehensive lesson plan with objectives, activities, and assessments..."):
                            stream = get_llm_response(tool, {"topic": topic, "duration": duration})
                            output = st.write_stream(stream)
                            st.session_state.last_tool_output = {'tool': tool, 'content': output}
                            add_history_entry("Lesson Plan", {"topic": topic, "duration": duration}, output)
                            just_generated = True
                    else:
                        st.warning("Please enter a topic or objective.")

            # --- TOOL: Report Card ---
            elif tool == "Report Card":
                st.title("📊 Report Card Comment Generator")
                st.write("Generate report card comments with a student's strengths and areas for growth.")
                name = st.text_input("Student's Name")
                strengths = st.text_area("Strengths", placeholder="e.g., Participates well in class, strong analytical skills")
                growth_areas = st.text_area("Areas for Growth", placeholder="e.g., Needs to show work in math, can be distracted")
                if st.button("Generate Comments", key="rc_gen"):
                    if name and strengths and growth_areas:
                        with st.spinner("📊 Crafting personalized report card comments highlighting strengths and growth opportunities..."):
                            stream = get_llm_response(tool, {"name": name, "strengths": strengths, "growth_areas": growth_areas})
                            output = st.write_stream(stream)
                            st.session_state.last_tool_output = {'tool': tool, 'content': output}
                            add_history_entry("Report Card", {"name": name, "strengths": strengths, "growth_areas": growth_areas}, output)
                            just_generated = True
                    else:
                        st.warning("Please fill in all fields.")

            # --- TOOL: Essay Grader ---
            elif tool == "Essay Grader":
                st.title("🎓 Essay Grader")
                st.write("Grade an essay based on a rubric or criteria.")
                essay = st.text_area("Paste the student's essay here:", height=300)
                rubric = st.text_area("Paste the grading rubric/criteria:", height=150, placeholder="e.g., Clarity: /5, Argument: /5, Grammar: /5")
                if st.button("Grade Essay", key="eg_gen"):
                    if essay and rubric:
                        with st.spinner("🎓 Analyzing essay structure, arguments, grammar, and coherence against the provided rubric..."):
                            stream = get_llm_response(tool, {"essay": essay, "rubric": rubric})
                            output = st.write_stream(stream)
                            st.session_state.last_tool_output = {'tool': tool, 'content': output}
                            add_history_entry("Essay Grader", {"essay": essay[:100] + "...", "rubric": rubric}, output)
                            just_generated = True
                    else:
                        st.warning("Please provide both the essay and the rubric.")

            # --- TOOL: PPT Generator ---
            elif tool == "PPT Generator":
                st.title("🖼️ PPT Generator")
                st.write("Generate a presentation outline based on a topic.")
                topic = st.text_input("Presentation Topic", placeholder="e.g., The History of Artificial Intelligence")
                num_slides = st.slider("Number of Slides", 5, 20, 10)
                if st.button("Generate PPT Outline", key="ppt_gen"):
                    if topic:
                        with st.spinner("🖼️ Designing engaging presentation slides with key points, visuals, and speaker notes..."):
                            stream = get_llm_response(tool, {"topic": topic, "num_slides": num_slides})
                            output = st.write_stream(stream)
                            st.session_state.last_tool_output = {'tool': tool, 'content': output}
                            add_history_entry("PPT Generator", {"topic": topic, "num_slides": num_slides}, output)
                            just_generated = True
                    else:
                        st.warning("Please enter a topic.")

            # --- TOOL: Question Paper Creation ---
            elif tool == "Question Paper Creation":
                st.title("📜 Question Paper Creator")
                st.write("Create a question paper for any topic, text, no. of questions and exam-type.")
                topic = st.text_input("Topic/Text", placeholder="e.g., World War II")
                num_q = st.number_input("Number of Questions", min_value=1, max_value=50, value=10)
                exam_type = st.selectbox("Exam Type", ["Semester", "Yearly", "Class Test"])
                if st.button("Create Question Paper", key="qp_gen"):
                    if topic:
                        with st.spinner("📜 Designing comprehensive question paper with varied question types and difficulty levels..."):
                            stream = get_llm_response(tool, {"topic": topic, "num_questions": num_q, "exam_type": exam_type})
                            output = st.write_stream(stream)
                            st.session_state.last_tool_output = {'tool': tool, 'content': output}
                            add_history_entry("Question Paper Creation", {"topic": topic, "num_questions": num_q, "exam_type": exam_type}, output)
                            just_generated = True
                    else:
                        st.warning("Please enter a topic.")

            # --- TOOL: Paper Correction ---
            elif tool == "Paper Correction":
                st.title("✅ Paper Correction Assistant")
                st.write("Generate a correction, summary, and marks based on the answer sheet provided.")
                answer_sheet = st.file_uploader("Upload Answer Sheet (Image or PDF)", type=['png', 'jpg', 'jpeg', 'pdf'])
                if st.button("Correct Paper", key="pc_gen"):
                    if answer_sheet:
                        with st.spinner("✅ Analyzing answer sheet for accuracy, completeness, and providing detailed feedback with marks..."):
                            stream = get_llm_response(tool, {"file_name": answer_sheet.name})
                            output = st.write_stream(stream)
                            st.session_state.last_tool_output = {'tool': tool, 'content': output}
                            add_history_entry("Paper Correction", {"file_name": answer_sheet.name}, output)
                            just_generated = True
                    else:
                        st.warning("Please upload an answer sheet.")

            # --- TOOL: YouTube or URL Summary ---
            elif tool == "YouTube or URL Summary":
                st.title("🔗 YouTube or URL Summary")
                st.write("Read or scrape a website/video and provide a summary with insights.")
                url = st.text_input("Enter YouTube Link or Website URL", placeholder="https://...")
                action_type = st.selectbox("Choose Action", ["Summary (100 words)", "Summary (500 words)", "Summary (1000 words)", "General Summary", "Complete Text", "Translation"])
                
                if st.button("Generate", key="url_sum_gen"):
                    if url:
                        with st.spinner("🔗 Fetching content from URL and analyzing for key insights and summaries..."):
                            content_text = fetch_url_content(url)
                            if "Error" in content_text:
                                st.error(content_text)
                            else:
                                stream = get_llm_response(tool, {"content": content_text, "action": action_type, "url": url})
                                output = st.write_stream(stream)
                                st.session_state.last_tool_output = {'tool': tool, 'content': output}
                                add_history_entry("YouTube or URL Summary", {"url": url, "action": action_type}, output)
                                just_generated = True
                    else:
                        st.warning("Please enter a URL.")

            # --- TOOL: Google Search Info ---
            elif tool == "Google Search Info":
                st.title("🔍 Google Search Info")
                st.write("Pull the latest information from Google on any topic.")
                topic = st.text_input("Enter Topic", placeholder="e.g., Latest advancements in AI")
                
                if st.button("Fetch Info", key="gs_info_gen"):
                    if topic:
                        with st.spinner("🔍 Searching Google and analyzing search results for the latest information..."):
                            search_results = fetch_google_search(topic)
                            if "Error" in search_results:
                                st.error(search_results)
                            else:
                                stream = get_llm_response(tool, {"topic": topic, "search_results": search_results})
                                output = st.write_stream(stream)
                                st.session_state.last_tool_output = {'tool': tool, 'content': output}
                                add_history_entry("Google Search Info", {"topic": topic}, output)
                                just_generated = True
                    else:
                        st.warning("Please enter a topic.")

            # --- TOOL: Google News ---
            elif tool == "Google News":
                st.title("📰 Google News")
                st.write("Get the latest news headlines and summaries from Google News.")
                topic = st.text_input("Enter News Topic", placeholder="e.g., Space Exploration")
                
                if st.button("Fetch News", key="gs_news_gen"):
                    if topic:
                        with st.spinner("📰 Fetching latest news headlines and summarizing key stories and events..."):
                            news_results = fetch_google_news(topic)
                            if "Error" in news_results:
                                st.error(news_results)
                            else:
                                stream = get_llm_response(tool, {"topic": topic, "news_results": news_results})
                                output = st.write_stream(stream)
                                st.session_state.last_tool_output = {'tool': tool, 'content': output}
                                add_history_entry("Google News", {"topic": topic}, output)
                                just_generated = True
                    else:
                        st.warning("Please enter a topic.")

            # --- CONSOLIDATED OUTPUT DISPLAY ---
            if st.session_state.get('last_tool_output') and st.session_state['last_tool_output']['tool'] == tool and not just_generated:
                st.success("Content Generated!")
                content = st.session_state['last_tool_output']['content']
                
                # Enhanced content display with better styling
                st.markdown('<div class="generated-content">', unsafe_allow_html=True)
                
                # Determine rendering method based on tool type
                markdown_tools = [
                    "Worksheet Generator", "Lesson Plan", "Essay Grader", 
                    "PPT Generator", "Question Paper Creation", "Paper Correction",
                    "YouTube or URL Summary", "Google Search Info", "Google News"
                ]
                if tool in markdown_tools:
                    st.markdown(content)
                else:
                    st.write(content)
                
                st.markdown('</div>', unsafe_allow_html=True)
                
                render_download_section(content, "final_output")

        else:
            # Show Grid of Tools
            st.title("✨ Magic Tools")
            st.write("Select a specialized AI service to get started")
            
            # Enhanced tools list with categories
            tools_list = [
                {
                    "name": "YouTube Generator", 
                    "icon": "🎬", 
                    "desc": "Generate guiding questions aligned to a YouTube video.",
                    "category": "Content Analysis",
                    "color": "#667eea"
                },
                {
                    "name": "Text Question", 
                    "icon": "❓", 
                    "desc": "Generate text-dependent questions for students based on any text.",
                    "category": "Assessment",
                    "color": "#28a745"
                },
                {
                    "name": "Worksheet Generator", 
                    "icon": "📝", 
                    "desc": "Generate a Worksheet based on any topic or text.",
                    "category": "Content Creation",
                    "color": "#17a2b8"
                },
                {
                    "name": "MCQ Generator", 
                    "icon": "📝", 
                    "desc": "Create a multiple choice assessment based on any topic.",
                    "category": "Assessment",
                    "color": "#28a745"
                },
                {
                    "name": "Text Summarizer", 
                    "icon": "✂️", 
                    "desc": "Take any text and summarize it in whatever length you choose.",
                    "category": "Content Analysis",
                    "color": "#667eea"
                },
                {
                    "name": "Text Rewriter", 
                    "icon": "✍️", 
                    "desc": "Take any text and rewrite it with custom criteria.",
                    "category": "Content Creation",
                    "color": "#17a2b8"
                },
                {
                    "name": "Proof Read", 
                    "icon": "🧐", 
                    "desc": "Have any text proofread for grammar, spelling, and clarity.",
                    "category": "Content Analysis",
                    "color": "#667eea"
                },
                {
                    "name": "Lesson Plan", 
                    "icon": "🍎", 
                    "desc": "Generate a lesson plan for a topic or objective.",
                    "category": "Planning",
                    "color": "#ffc107"
                },
                {
                    "name": "Report Card", 
                    "icon": "📊", 
                    "desc": "Generate report card comments based on strengths/growth.",
                    "category": "Assessment",
                    "color": "#28a745"
                },
                {
                    "name": "Essay Grader", 
                    "icon": "🎓", 
                    "desc": "Grade an essay based on a rubric or criteria.",
                    "category": "Assessment",
                    "color": "#28a745"
                },
                {
                    "name": "PPT Generator", 
                    "icon": "🖼️", 
                    "desc": "Generate a presentation outline based on a topic.",
                    "category": "Content Creation",
                    "color": "#17a2b8"
                },
                {
                    "name": "Question Paper Creation", 
                    "icon": "📜", 
                    "desc": "Create a question paper for any topic and exam-type.",
                    "category": "Assessment",
                    "color": "#28a745"
                },
                {
                    "name": "Paper Correction", 
                    "icon": "✅", 
                    "desc": "Generate a correction and marks based on an answer sheet.",
                    "category": "Assessment",
                    "color": "#28a745"
                },
                {
                    "name": "YouTube or URL Summary", 
                    "icon": "🔗", 
                    "desc": "Summarize a YouTube video or Website URL with insights.",
                    "category": "Content Analysis",
                    "color": "#667eea"
                },
                {
                    "name": "Google Search Info", 
                    "icon": "🔍", 
                    "desc": "Pull latest information from Google on any topic.",
                    "category": "Research",
                    "color": "#6f42c1"
                },
                {
                    "name": "Google News", 
                    "icon": "📰", 
                    "desc": "Get latest news headlines from Google News.",
                    "category": "Research",
                    "color": "#6f42c1"
                }
            ]
            
            # Grid Layout: Iterate in chunks of 3 to ensure proper alignment
            for i in range(0, len(tools_list), 3):
                cols = st.columns(3)
                for j in range(3):
                    if i + j < len(tools_list):
                        tool = tools_list[i + j]
                        with cols[j]:
                            # Enhanced card layout with category tag
                            st.markdown(f"""
                            <div class="service-card" data-tool-name="{tool['name']}" style="border-left: 4px solid {tool['color']}">
                                <div class="service-card-content">
                                    <div class="category-tag">{tool['category']}</div>
                                    <div class="service-icon">{tool['icon']}</div>
                                    <div class="service-title">{tool['name']}</div>
                                    <div class="service-desc">{tool['desc']}</div>
                                </div>
                            </div>
                            """, unsafe_allow_html=True)
                            
                            # Form submission for tool selection
                            with st.form(key=f"form_{i+j}"):
                                submitted = st.form_submit_button("Select", type="primary", use_container_width=True, help=f"Click to open {tool['name']}")
                                if submitted:
                                    st.session_state.current_tool = tool['name']
                                    st.rerun()
